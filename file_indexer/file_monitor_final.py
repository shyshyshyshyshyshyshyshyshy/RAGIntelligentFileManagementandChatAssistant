# file_monitor_dify_chatflow.py - 使用Dify Chatflow解析文件
import os
import sys
import time
import json
import hashlib
import uuid
from datetime import datetime
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from dotenv import load_dotenv
import requests
import logging
import re
import pandas as pd
from pptx import Presentation
import PyPDF2
import docx
import time
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from flask import send_from_directory  # 新增这一行
import threading  # 新增这一行
import webbrowser  # 新增这一行

# 加载环境变量
load_dotenv()

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('file_monitor.log', encoding='utf-8'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)

try:
    from PIL import Image, ExifTags
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logger.warning("PIL/Pillow未安装，图片分析功能将受限")

class FileConverter:
    """文件转换器"""
    
    @staticmethod
    def convert_doc_to_docx(file_path):
        """将.doc文件转换为.docx格式"""
        try:
            file_name = os.path.basename(file_path)
            logger.info(f"开始转换.doc文件为.docx: {file_name}")
            
            # 创建临时目录
            import tempfile
            temp_dir = tempfile.mkdtemp()
            base_name = os.path.splitext(file_name)[0]
            docx_file_path = os.path.join(temp_dir, f"{base_name}_converted.docx")
            
            # 使用改进的转换方法
            try:
                return FileConverter._convert_with_improved_method(file_path, docx_file_path)
            except Exception as e:
                logger.warning(f"改进转换方法失败: {str(e)}")
                return None
            
        except Exception as e:
            logger.error(f"文件转换异常: {str(e)}")
            return None
    
    @staticmethod
    def _convert_with_improved_method(doc_path, docx_path):
        """使用改进的方法转换.doc文件"""
        try:
            # 方法1: 尝试使用外部工具
            try:
                return FileConverter._convert_with_external_tools(doc_path, docx_path)
            except Exception as e:
                logger.debug(f"外部工具转换失败: {str(e)}")
            
            # 方法2: 使用二进制分析提取文本
            try:
                return FileConverter._convert_with_binary_analysis(doc_path, docx_path)
            except Exception as e:
                logger.debug(f"二进制分析转换失败: {str(e)}")
            
            # 方法3: 创建简单的占位文档
            return FileConverter._create_placeholder_docx(doc_path, docx_path)
            
        except Exception as e:
            raise Exception(f"所有转换方法都失败: {str(e)}")
    
    @staticmethod
    def _convert_with_external_tools(doc_path, docx_path):
        """使用外部工具转换.doc文件"""
        try:
            # 检查是否安装了外部工具
            import subprocess
            
            # 尝试使用antiword
            try:
                result = subprocess.run(['antiword', doc_path], capture_output=True, text=True, timeout=30, encoding='utf-8', errors='ignore')
                if result.returncode == 0 and result.stdout.strip():
                    content = result.stdout
                    logger.info("使用antiword成功提取文本内容")
                    return FileConverter._create_docx_from_text(content, docx_path, os.path.basename(doc_path))
            except Exception as e:
                logger.debug(f"antiword失败: {str(e)}")
            
            # 尝试使用catdoc
            try:
                result = subprocess.run(['catdoc', doc_path], capture_output=True, text=True, timeout=30, encoding='utf-8', errors='ignore')
                if result.returncode == 0 and result.stdout.strip():
                    content = result.stdout
                    logger.info("使用catdoc成功提取文本内容")
                    return FileConverter._create_docx_from_text(content, docx_path, os.path.basename(doc_path))
            except Exception as e:
                logger.debug(f"catdoc失败: {str(e)}")
            
            raise Exception("外部工具不可用或转换失败")
            
        except Exception as e:
            raise Exception(f"外部工具转换错误: {str(e)}")
    
    @staticmethod
    def _convert_with_binary_analysis(doc_path, docx_path):
        """通过二进制分析提取文本内容"""
        try:
            with open(doc_path, 'rb') as f:
                content = f.read()
            
            # 尝试多种编码提取文本
            text_content = ""
            
            # 尝试UTF-8编码
            try:
                decoded = content.decode('utf-8', errors='ignore')
                # 提取可打印字符
                printable_chars = ''.join(char for char in decoded if char.isprintable() or char in '\n\r\t')
                if len(printable_chars) > 100:  # 确保有足够的内容
                    text_content = printable_chars
                    logger.info("使用UTF-8编码提取文本成功")
            except:
                pass
            
            # 如果UTF-8失败，尝试其他编码
            if not text_content:
                encodings = ['gbk', 'gb2312', 'latin-1', 'cp1252']
                for encoding in encodings:
                    try:
                        decoded = content.decode(encoding, errors='ignore')
                        printable_chars = ''.join(char for char in decoded if char.isprintable() or char in '\n\r\t')
                        if len(printable_chars) > 100:
                            text_content = printable_chars
                            logger.info(f"使用{encoding}编码提取文本成功")
                            break
                    except:
                        continue
            
            if text_content:
                # 清理文本内容
                cleaned_content = FileConverter._clean_extracted_text(text_content)
                return FileConverter._create_docx_from_text(cleaned_content, docx_path, os.path.basename(doc_path))
            else:
                raise Exception("无法从二进制数据中提取文本")
                
        except Exception as e:
            raise Exception(f"二进制分析失败: {str(e)}")
    
    @staticmethod
    def _clean_extracted_text(text):
        """清理提取的文本内容"""
        if not text:
            return ""
        
        # 移除控制字符
        import re
        cleaned = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', text)
        
        # 移除过多的空白字符
        cleaned = re.sub(r'\n\s*\n', '\n\n', cleaned)  # 保留段落间隔
        cleaned = re.sub(r'[ \t]+', ' ', cleaned)  # 压缩空格
        
        # 限制文本长度
        if len(cleaned) > 10000:
            cleaned = cleaned[:10000] + "\n\n[内容已截断]"
        
        return cleaned.strip()
    
    @staticmethod
    def _create_placeholder_docx(doc_path, docx_path):
        """创建占位符文档"""
        try:
            from docx import Document
            doc = Document()
            
            file_name = os.path.basename(doc_path)
            doc.add_heading(f"文档转换自: {file_name}", 0)
            doc.add_paragraph("注意: 这是一个自动生成的占位符文档。")
            doc.add_paragraph("原始.doc文件无法被完全解析，可能包含二进制数据或特殊格式。")
            doc.add_paragraph(f"文件大小: {os.path.getsize(doc_path)} 字节")
            doc.add_paragraph(f"转换时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            
            doc.save(docx_path)
            logger.info(f"创建占位符文档: {file_name}")
            return docx_path
            
        except Exception as e:
            raise Exception(f"创建占位符文档失败: {str(e)}")
    
    @staticmethod
    def _create_docx_from_text(content, docx_path, original_filename):
        """从文本内容创建.docx文件"""
        try:
            from docx import Document
            doc = Document()
            
            # 添加标题
            doc.add_heading(f"转换自: {original_filename}", 0)
            
            # 清理内容，确保XML兼容
            cleaned_content = FileConverter._clean_text_for_xml(content)
            
            if not cleaned_content.strip():
                doc.add_paragraph("警告: 提取的内容为空或无法解析")
                doc.add_paragraph("原始文件可能包含加密内容、损坏数据或特殊格式。")
            else:
                # 添加内容
                paragraphs = cleaned_content.split('\n')
                for para in paragraphs:
                    if para.strip():
                        try:
                            # 确保段落文本是XML兼容的
                            safe_para = FileConverter._make_xml_safe(para)
                            doc.add_paragraph(safe_para)
                        except Exception as e:
                            # 如果添加段落失败，记录并跳过
                            logger.debug(f"跳过不兼容XML的段落: {str(e)}")
                            continue
            
            # 添加文件信息
            doc.add_heading("文件信息", 1)
            doc.add_paragraph(f"原始文件名: {original_filename}")
            doc.add_paragraph(f"转换时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            doc.add_paragraph(f"内容长度: {len(cleaned_content)} 字符")
            
            # 保存.docx文件
            doc.save(docx_path)
            logger.info(f"成功创建.docx文件: {original_filename}")
            return docx_path
            
        except Exception as e:
            raise Exception(f"创建.docx文件错误: {str(e)}")
    
    @staticmethod
    def _make_xml_safe(text):
        """确保文本是XML安全的"""
        if not text:
            return ""
        
        # 替换XML特殊字符
        safe_text = text.replace('&', '&amp;')
        safe_text = safe_text.replace('<', '&lt;')
        safe_text = safe_text.replace('>', '&gt;')
        safe_text = safe_text.replace('"', '&quot;')
        safe_text = safe_text.replace("'", '&apos;')
        
        # 移除控制字符
        import re
        safe_text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', safe_text)
        
        return safe_text
    
    @staticmethod
    def _clean_text_for_xml(text):
        """清理文本，确保XML兼容"""
        if not text:
            return ""
        
        # 移除控制字符和NULL字节
        import re
        # 移除ASCII控制字符（除了换行符和制表符）
        cleaned = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F]', '', text)
        # 移除Unicode控制字符
        cleaned = re.sub(r'[\u2000-\u200F\u2028-\u202F\u205F-\u206F]', '', cleaned)
        
        return cleaned
    
class Config:
    """配置"""
    DIFY_BASE_URL = os.getenv('DIFY_BASE_URL', 'http://localhost').rstrip('/')
    DATASET_API_KEY = os.getenv('DIFY_API_KEY', 'dataset-zqGccO9VowfmI7bPG6opOh5C')
    
    
    
    # 添加图片相关配置
    MAX_IMAGE_SIZE = int(os.getenv('MAX_IMAGE_SIZE', '10485760'))  # 10MB
    IMAGE_UPLOAD_ENABLED = os.getenv('IMAGE_UPLOAD_ENABLED', 'false').lower() == 'true'
    
    # 两个知识库ID
    TXT_KNOWLEDGE_BASE_ID = os.getenv('DIFY_KNOWLEDGE_BASE_ID', '1f0cc924-cba1-4113-83eb-dca99b0a31f9')
    ORIGINAL_KNOWLEDGE_BASE_ID = os.getenv('ORIGINAL_KNOWLEDGE_BASE_ID', 'b59ade74-3ce2-4857-8d67-8f3f0faedab2')
    
    #父子模式配置
    PARENT_CHILD_KB_ID = os.getenv('PARENT_CHILD_KB_ID', '1388750e-551b-4084-b699-17091a5b8364')  # 您新创建的父子模式知识库ID
    PARENT_CHILD_KB_ENABLED = os.getenv('PARENT_CHILD_KB_ENABLED', 'true').lower() == 'true'

    # Chatflow配置
    CHATFLOW_APP_ID = os.getenv('WORKFLOW_APP_ID', 'app-A9xu8Jb4GeGZAOvHhv3capwV')
    CHATFLOW_API_KEY = os.getenv('WORKFLOW_API_KEY', 'app-A9xu8Jb4GeGZAOvHhv3capwV')
    
    MONITOR_DIR = os.getenv('MONITOR_DIR', 'D:/code/python/ALLOWED_FILE_DIR')
    ALLOWED_EXTENSIONS = tuple(os.getenv('ALLOWED_EXTENSIONS', 
        '.txt,.docx,.doc,.pdf,.pptx,.xlsx,.csv,.md,.jpg,.jpeg,.png,.gif,.bmp,.webp,.tiff').split(','))
    PROCESS_INTERVAL = int(os.getenv('PROCESS_INTERVAL', '5'))
    API_TIMEOUT = int(os.getenv('API_TIMEOUT', '60'))

    @property
    def ACTUAL_ORIGINAL_KB_ID(self):
        """动态决定使用哪个原文件知识库"""
        if self.PARENT_CHILD_KB_ENABLED and self.PARENT_CHILD_KB_ID:
            return self.PARENT_CHILD_KB_ID
        else:
            return self.ORIGINAL_KNOWLEDGE_BASE_ID

config = Config()

# 修复 DifyChatflowAnalyzer 类，添加 extract_text_content 方法
# 修改 DifyChatflowAnalyzer 类中的 _process_ai_response 方法

class DifyChatflowAnalyzer:
    """Dify Chatflow分析器"""
    
    def __init__(self):
        # 初始化 SimpleFileAnalyzer 实例用于提取文本内容
        self.simple_analyzer = SimpleFileAnalyzer()
    
    def analyze_with_chatflow(self, file_path):
        """使用Dify Chatflow分析文件 - 支持图片"""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            is_image = file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff']
        
            if is_image:
                return self._analyze_image_with_chatflow(file_path)
            else:
                return self._analyze_document_with_chatflow(file_path)
    
        except FileNotFoundError:
            logger.error(f"文件不存在: {file_path}")
            return None
        except PermissionError:
            logger.error(f"无权限访问文件: {file_path}")
            return None
        except Exception as e:
            logger.error(f"Chatflow分析未知异常: {str(e)}")
            return None

    def _analyze_image_with_chatflow(self, file_path):
        """专门分析图片的Chatflow方法"""
        try:
            file_name = os.path.basename(file_path)
            logger.info(f"开始使用Dify Chatflow分析图片: {file_name}")
            
            # 检查PIL是否可用
            if not PIL_AVAILABLE:
                logger.warning("PIL/Pillow未安装，无法分析图片内容")
                return None
            
            # 提取图片信息
            image_info_str = EnhancedFileAnalyzer._extract_image_content(file_path)
            if not image_info_str or "图片处理错误" in image_info_str:
                logger.warning(f"图片信息提取失败: {file_name}")
                return None
            
            # 构造图片分析专用提示词
            url = f"{config.DIFY_BASE_URL}/v1/chat-messages"
            headers = {
                "Authorization": f"Bearer {config.CHATFLOW_API_KEY}",
                "Content-Type": "application/json"
            }
            
            prompt = f"""
请分析这张图片，用3-4句话描述图片内容，包括：
1. 主要场景和背景
2. 图片中的主体对象  
3. 颜色、光线等视觉特征
4. 整体氛围或情感

图片信息:
文件名: {file_name}

请直接输出描述内容，不要添加思考过程或分析步骤。
            """
            
            data = {
                "inputs": {},
                "query": prompt,
                "response_mode": "blocking",
                "user": f"file_monitor_{hashlib.md5(file_name.encode()).hexdigest()[:8]}"
            }
            
            response = requests.post(url, headers=headers, json=data, timeout=config.API_TIMEOUT)
            
            if response.status_code == 200:
                result = response.json()
                logger.info(f"Dify Chatflow图片分析成功: {file_name}")
                
                # 处理返回结果
                processed_result = self._process_image_ai_response(result, file_name)
                return processed_result
            else:
                logger.error(f"Dify Chatflow图片分析失败: {response.status_code}")
                return None
                
        except Exception as e:
            logger.error(f"图片Chatflow分析异常: {str(e)}")
            return None
    
    def _analyze_document_with_chatflow(self, file_path):
        """原有的文档分析方法"""
        try:
            file_name = os.path.basename(file_path)
            logger.info(f"开始使用Dify Chatflow分析文档: {file_name}")
            
            # 使用 SimpleFileAnalyzer 提取文本内容
            file_content = EnhancedFileAnalyzer.extract_text_content(file_path)
            if not file_content or len(file_content.strip()) == 0:
                logger.warning(f"文件内容为空，跳过Chatflow分析: {file_name}")
                return None
            
            # 构造Chatflow请求
            url = f"{config.DIFY_BASE_URL}/v1/chat-messages"
            
            headers = {
                "Authorization": f"Bearer {config.CHATFLOW_API_KEY}",
                "Content-Type": "application/json"
            }
            
            # 更严格的提示词，强制大模型不输出思考过程
            prompt = f"""
请严格按照以下格式分析文档，不要包含任何思考过程、分析步骤或解释：

文件类型: [作业/实验报告/学术论文/项目报告/技术文档/学习笔记等]
内容总结: [首先推断这是什么类型的文件：是作业还是实验报告文档还是论文还是什么，可以有多个推断。然后用3-4句话简洁总结文档关键内容]

文档信息:
文件名: {file_name}
内容: {file_content[:4000]}

请直接输出格式化的结果，不要添加任何其他内容。
            """
            
            data = {
                "inputs": {},
                "query": prompt,
                "response_mode": "blocking",
                "user": f"file_monitor_{hashlib.md5(file_name.encode()).hexdigest()[:8]}"
            }
            
            logger.debug(f"Chatflow请求URL: {url}")
            
            response = requests.post(
                url,
                headers=headers,
                json=data,
                timeout=config.API_TIMEOUT
            )
            
            logger.debug(f"Chatflow响应状态码: {response.status_code}")
            
            if response.status_code == 200:
                result = response.json()
                logger.info(f"Dify Chatflow分析成功: {file_name}")
                
                # 对返回结果进行后处理，确保格式正确
                processed_result = self._process_ai_response(result)
                return processed_result
            else:
                logger.error(f"Dify Chatflow分析失败: {response.status_code} - {response.text}")
                return None
                
        except requests.exceptions.Timeout:
            logger.error("Dify Chatflow请求超时")
            return None
        except requests.exceptions.ConnectionError:
            logger.error("Dify Chatflow连接错误，请检查网络连接和DIFY_BASE_URL")
            return None
        except Exception as e:
            logger.error(f"Dify Chatflow分析异常: {str(e)}")
            return None
    
    def _process_image_ai_response(self, result, file_name):
        """处理图片分析的AI返回结果"""
        try:
            answer = result.get('answer', '')
            logger.debug(f"原始图片AI回答: {answer}")
            
            # 清理思考标签
            cleaned_answer = self._remove_thought_tags(answer)
            
            # 构建标准格式的返回结果
            if cleaned_answer.strip():
                # 对于图片，文件类型固定为"图片"
                processed_answer = f"文件类型: 图片\n内容总结: {cleaned_answer}"
                result['answer'] = processed_answer
            else:
                result['answer'] = "文件类型: 图片\n内容总结: 无法分析图片内容"
            
            return result
        except Exception as e:
            logger.error(f"处理图片AI响应失败: {str(e)}")
            return result
    
    def _process_ai_response(self, result):
        """处理AI返回结果，确保格式正确"""
        try:
            answer = result.get('answer', '')
            logger.debug(f"原始AI回答: {answer}")
            
            # 过滤掉思考标签内的内容
            cleaned_answer = self._remove_thought_tags(answer)
            logger.debug(f"过滤思考标签后: {cleaned_answer}")
            
            # 如果AI没有按照格式返回，尝试提取关键信息
            if '文件类型:' not in cleaned_answer or '内容总结:' not in cleaned_answer:
                logger.warning("AI返回格式不符合要求，尝试提取关键信息")
                
                # 提取文件类型
                doc_type = self._extract_doc_type(cleaned_answer)
                # 提取内容总结
                summary = self._extract_summary(cleaned_answer)
                
                # 重新构建标准格式
                processed_answer = f"文件类型: {doc_type}\n内容总结: {summary}"
                result['answer'] = processed_answer
            else:
                result['answer'] = cleaned_answer
            
            return result
        except Exception as e:
            logger.error(f"处理AI响应失败: {str(e)}")
            return result
    
    def _remove_thought_tags(self, text):
        """移除思考标签内的内容"""
        try:
            # 定义思考标签的正则表达式
            thought_pattern = r'<think>.*?</think>'
            
            # 移除所有思考标签及其内容
            cleaned_text = re.sub(thought_pattern, '', text, flags=re.DOTALL)
            
            # 如果移除后文本为空或只有空白，返回原始文本
            if not cleaned_text.strip():
                return text
                
            return cleaned_text.strip()
        except Exception as e:
            logger.warning(f"移除思考标签失败: {str(e)}")
            return text
    
    def _extract_doc_type(self, answer):
        """从AI回答中提取文件类型"""
        # 常见文档类型关键词
        doc_types = {
            '作业': '学生作业',
            'assignment': '学生作业',
            'homework': '学生作业',
            '实验': '实验报告', 
            'experiment': '实验报告',
            'lab': '实验报告',
            '论文': '学术论文',
            'thesis': '学术论文',
            'paper': '学术论文',
            '报告': '项目报告',
            'report': '项目报告',
            '设计': '设计文档',
            'design': '设计文档',
            '笔记': '学习笔记',
            'note': '学习笔记',
            '技术': '技术文档',
            '开发': '技术文档',
            '代码': '技术文档'
        }
        
        answer_lower = answer.lower()
        for keyword, doc_type in doc_types.items():
            if keyword in answer_lower:
                return doc_type
        
        return "通用文档"
    
    def _extract_summary(self, answer):
        """从AI回答中提取内容总结"""
        # 移除思考过程标记
        summary = answer.replace('首先，', '').replace('接下来，', '').replace('然后，', '')
        summary = summary.replace('嗯，', '').replace('首先我需要', '').replace('我需要', '')
        
        # 提取第一句话或前100个字符
        sentences = re.split(r'[。！？]', summary)
        if sentences:
            first_sentence = sentences[0].strip()
            if len(first_sentence) > 10:  # 确保不是空句子
                if len(first_sentence) > 100:
                    return first_sentence[:100] + "..."
                return first_sentence
        
        # 如果无法提取，返回简化的总结
        if len(answer) > 100:
            return answer[:100] + "..."
        return answer if answer else "无法生成内容总结"
    def check_multimodal_support(self):
        """检查Dify是否支持多模态分析"""
        try:
            # 测试端点 - 使用模型列表接口
            test_url = f"{config.DIFY_BASE_URL}/v1/models"
            headers = {"Authorization": f"Bearer {config.CHATFLOW_API_KEY}"}
            
            response = requests.get(test_url, headers=headers, timeout=10)
            if response.status_code == 200:
                result = response.json()
                print("✅ Dify API连接成功")
                print(f"可用模型: {result}")
                return True
            else:
                print(f"❌ Dify API连接失败: {response.status_code}")
                return False
        except Exception as e:
            print(f"❌ 测试失败: {str(e)}")
            return False
# 修改 DifyChatflowIndexGenerator 类中的 _parse_ai_response 方法

class DifyChatflowIndexGenerator:
    """Dify Chatflow索引生成器"""
    
    def __init__(self):
        self.chatflow_analyzer = DifyChatflowAnalyzer()
        self.fallback_analyzer = SimpleFileAnalyzer()
        self.info_extractor = FileInfoExtractor()
    
    def _enhance_summary_with_filename(self, content_summary, original_filename, max_length=250):
        """在内容总结中增强文件名信息"""
        # 基础增强：在总结开头添加文件名
        enhanced = f"文件【{original_filename}】相关内容：{content_summary}"
    
        # 处理长度限制：优先保留文件名，适当截断内容
        if len(enhanced) > max_length:
                # 计算文件名部分占用的长度
            filename_part = f"文件【{original_filename}】相关内容："
            available_space = max_length - len(filename_part)
        
            if available_space > 10:  # 确保有最小内容空间
                # 截断内容总结但保留文件名
                truncated_content = content_summary[:available_space] + "..."
                enhanced = filename_part + truncated_content
            else:
                # 极端情况：只保留关键信息
                enhanced = f"文件【{original_filename}】相关文档"
    
        return enhanced
    
    def generate_index_file(self, file_path):
        """生成索引文件（优先使用Dify Chatflow）"""
        try:
            file_name = os.path.basename(file_path)
            
            # 跳过索引文件
            if any(marker in file_name for marker in ['_summary', '_index', '_workflow', '_dify', '_chatflow']):
                return None
            
            logger.info(f"开始处理文件: {file_name}")
            
            # 获取文件信息
            file_info = self.info_extractor.extract_file_info(file_path)
            if not file_info:
                return None
            
            # 1. 优先使用Dify Chatflow分析
            chatflow_result = self.chatflow_analyzer.analyze_with_chatflow(file_path)
            
            if chatflow_result:
                # Chatflow分析成功
                index_content = self._format_chatflow_index(file_info, chatflow_result, file_path)
                analysis_method = "Dify Chatflow分析"
                is_fallback = False
            else:
                # Chatflow失败，使用备用方案
                logger.info(f"Chatflow分析失败，使用备用方案: {file_name}")
                content = self.fallback_analyzer.extract_text_content(file_path)
                file_type = self.fallback_analyzer.infer_file_type(file_name, content)
                content_summary = self._simplify_content_summary(content)
                index_content = self._format_fallback_index(file_info, file_type, content_summary)
                analysis_method = "本地规则推断"
                is_fallback = True
            
            # 保存索引文件
            base_name = os.path.splitext(file_name)[0]
            suffix = "_fallback_index" if is_fallback else "_chatflow_index"
            index_filename = f"{base_name}{suffix}.txt"
            index_path = os.path.join(config.MONITOR_DIR, index_filename)
            
            with open(index_path, 'w', encoding='utf-8') as f:
                f.write(index_content)
            
            logger.info(f"索引文件已生成 ({analysis_method}): {index_filename}")
            return index_path
            
        except Exception as e:
            logger.error(f"索引文件生成失败: {file_path} - {str(e)}")
            return None
    
    def _format_chatflow_index(self, file_info, chatflow_result, file_path):
        """格式化Chatflow分析结果索引文件（严格精简版）"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # 提取Chatflow返回的分析结果
        analysis_result = chatflow_result.get('answer', '无分析结果')
        logger.debug(f"索引生成器收到的分析结果: {analysis_result}")
        
        # 解析AI返回的结果
        doc_type, content_summary = self._parse_ai_response(analysis_result, file_info['name'])
        enhanced_summary = self._enhance_summary_with_filename(content_summary, file_info['name'])
        logger.debug(f"解析后的文件类型: {doc_type}, 内容总结: {content_summary}")
        
        # 严格精简的索引内容
        index_content = f"""文件名: {file_info['name']}
文件路径: {file_info['path']}
修改时间: {file_info['update_time']}
文件类型: {doc_type}
内容总结: {enhanced_summary}"""

        return index_content
    
    def _parse_ai_response(self, analysis_result, file_name):
        """解析AI返回的结果"""
        try:
            logger.debug(f"开始解析AI响应: {analysis_result}")
            
            # 再次确保移除思考标签（双重保险）
            cleaned_result = self._remove_thought_tags(analysis_result)
            logger.debug(f"清理思考标签后: {cleaned_result}")
            
            # 尝试从标准格式中提取信息
            if '文件类型:' in cleaned_result and '内容总结:' in cleaned_result:
                lines = cleaned_result.split('\n')
                doc_type = "通用文档"
                content_summary = "无法生成内容总结"
                
                for line in lines:
                    line_stripped = line.strip()
                    if line_stripped.startswith('文件类型:'):
                        doc_type = line_stripped.replace('文件类型:', '').strip()
                        logger.debug(f"提取到文件类型: {doc_type}")
                    elif line_stripped.startswith('内容总结:'):
                        content_summary = line_stripped.replace('内容总结:', '').strip()
                        logger.debug(f"提取到内容总结: {content_summary}")
                
                # 如果内容总结包含思考过程，进一步清理
                if '首先' in content_summary or '接下来' in content_summary or '然后' in content_summary:
                    # 提取第一个句子的内容
                    sentences = re.split(r'[。！？]', content_summary)
                    if sentences:
                        content_summary = sentences[0].strip()
                
                return doc_type, content_summary
            else:
                # 如果格式不符合，使用备用推断
                logger.warning("AI响应格式不符合标准，使用备用方案")
                doc_type = self.fallback_analyzer.infer_file_type(file_name, cleaned_result)
                content_summary = self._simplify_content_summary(cleaned_result)
                return doc_type, content_summary
                
        except Exception as e:
            logger.warning(f"解析AI响应失败: {str(e)}")
            # 使用备用方案
            doc_type = self.fallback_analyzer.infer_file_type(file_name, "")
            content_summary = "AI分析结果格式异常"
            return doc_type, content_summary
    
    def _remove_thought_tags(self, text):
        """移除思考标签内的内容（双重保险）"""
        try:
            # 定义思考标签的正则表达式
            thought_pattern = r'<think>.*?</think>'
            
            # 移除所有思考标签及其内容
            cleaned_text = re.sub(thought_pattern, '', text, flags=re.DOTALL)
            
            # 如果移除后文本为空或只有空白，返回原始文本
            if not cleaned_text.strip():
                return text
                
            return cleaned_text.strip()
        except Exception as e:
            logger.warning(f"移除思考标签失败: {str(e)}")
            return text
    
    def _format_fallback_index(self, file_info, file_type, content_summary):
        """格式化备用方案索引文件（精简版）"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        enhanced_summary = self._enhance_summary_with_filename(content_summary, file_info['name'])
        # 精简的索引内容
        index_content = f"""文件索引信息
====================
文件名: {file_info['name']}
文件路径: {file_info['path']}
创建时间: {file_info['create_time']}
修改时间: {file_info['update_time']}
文件类型: {file_type}

内容总结: {enhanced_summary}

生成时间: {timestamp}
分析引擎: 本地规则推断
===================="""

        return index_content
    
    def _simplify_content_summary(self, content):
        """简化内容总结"""
        if not content or len(content.strip()) == 0:
            return "文件内容为空或无法读取"
        
        # 移除思考过程标记
        cleaned_content = re.sub(r'首先，|接下来，|然后，|嗯，|首先我需要|我需要', '', content)
        
        # 提取第一句话或前80个字符
        sentences = re.split(r'[。！？]', cleaned_content)
        if sentences:
            first_sentence = sentences[0].strip()
            if len(first_sentence) > 10:  # 确保不是空句子
                if len(first_sentence) > 80:
                    return first_sentence[:80] + "..."
                return first_sentence
        
        # 如果无法提取，返回简化的总结
        if len(content) > 80:
            return content[:80] + "..."
        return content
    
class SimpleFileAnalyzer:
    """基础文件分析器（用于备用方案）"""
    
    @staticmethod
    def extract_text_content(file_path):
        """提取文本内容（兼容EnhancedFileAnalyzer）"""
        return EnhancedFileAnalyzer.extract_text_content(file_path)
    
    @staticmethod
    def infer_file_type(file_name, content):
        """基于文件名和内容推断文件类型"""
        try:
            file_name_lower = file_name.lower()
            content_lower = content.lower() if content else ""
            
            # 基于文件名关键词匹配
            if any(keyword in file_name_lower for keyword in ['作业', 'assignment', 'homework']):
                return "学生作业"
            elif any(keyword in file_name_lower for keyword in ['实验', 'experiment', 'lab']):
                return "实验报告"
            elif any(keyword in file_name_lower for keyword in ['论文', 'thesis', 'paper']):
                return "学术论文"
            elif any(keyword in file_name_lower for keyword in ['报告', 'report']):
                return "项目报告"
            elif any(keyword in file_name_lower for keyword in ['设计', 'design']):
                return "设计文档"
            elif any(keyword in file_name_lower for keyword in ['笔记', 'note']):
                return "学习笔记"
            elif any(keyword in file_name_lower for keyword in ['技术', '开发', '代码']):
                return "技术文档"
            
            # 基于内容关键词匹配
            if any(keyword in content_lower for keyword in ['实验目的', '实验步骤', '实验结果']):
                return "实验报告"
            elif any(keyword in content_lower for keyword in ['摘要', '关键词', '参考文献']):
                return "学术论文"
            elif any(keyword in content_lower for keyword in ['需求分析', '设计思路', '实现方案']):
                return "项目报告"
            else:
                return "通用文档"
        except Exception as e:
            logger.warning(f"推断文件类型失败: {str(e)}")
            return "通用文档"

class EnhancedFileAnalyzer:
    """增强的文件分析器"""
    
    @staticmethod
    def extract_text_content(file_path):
        """增强的文本内容提取"""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
                    
            elif file_ext == '.docx':
                try:
                    import docx
                    doc = docx.Document(file_path)
                    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
                    return '\n'.join(paragraphs)
                except Exception as e:
                    return f"Word文档读取错误: {str(e)}"
                    
            elif file_ext == '.doc':
                # 旧版Word文档处理
                return EnhancedFileAnalyzer._extract_doc_content(file_path)
                    
            elif file_ext == '.pdf':
                return EnhancedFileAnalyzer._extract_pdf_content(file_path)
                    
            elif file_ext == '.xlsx':
                return EnhancedFileAnalyzer._extract_excel_content(file_path)
                    
            elif file_ext == '.pptx':
                return EnhancedFileAnalyzer._extract_ppt_content(file_path)
            
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff']:
                return EnhancedFileAnalyzer._extract_image_content(file_path)
            
            else:
                # 其他文本文件尝试读取
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read()
                except:
                    return ""
                    
        except Exception as e:
            return f"内容提取失败: {str(e)}"
        
    @staticmethod
    def _extract_image_content(file_path):
        """提取图片内容信息"""
        try:
            # 使用PIL获取图片基本信息
            from PIL import Image, ExifTags
            import base64
            
            with Image.open(file_path) as img:
                # 获取图片基本信息
                info = {
                    'format': img.format,
                    'size': img.size,
                    'mode': img.mode,
                    'filename': os.path.basename(file_path)
                }
                
                # 获取EXIF信息
                try:
                    exif_data = img._getexif()
                    if exif_data:
                        for tag, value in exif_data.items():
                            tag_name = ExifTags.TAGS.get(tag, tag)
                            info[f'exif_{tag_name}'] = str(value)
                except:
                    pass
                
                # 转换为base64用于大模型分析
                with open(file_path, 'rb') as f:
                    image_base64 = base64.b64encode(f.read()).decode('utf-8')
                    info['base64_data'] = image_base64
                
                return json.dumps(info, ensure_ascii=False)
                
        except Exception as e:
            return f"图片处理错误: {str(e)}"
        

    @staticmethod
    def _extract_doc_content_enhanced(file_path):
        """增强的.doc文件内容提取"""
        try:
            # 先尝试转换后提取
            converter = FileConverter()
            converted_path = converter.convert_doc_to_docx(file_path)
            
            if converted_path and os.path.exists(converted_path):
                try:
                    # 从转换后的.docx文件中提取内容
                    from docx import Document
                    doc = Document(converted_path)
                    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
                    content = '\n'.join(paragraphs)
                    
                    if content and len(content.strip()) > 0:
                        logger.info(f"从转换文件成功提取内容: {os.path.basename(file_path)}")
                        return content
                    else:
                        logger.warning(f"转换文件内容为空: {os.path.basename(file_path)}")
                except Exception as e:
                    logger.warning(f"从转换文件提取内容失败: {str(e)}")
                finally:
                    # 清理临时文件
                    if os.path.exists(converted_path):
                        try:
                            os.remove(converted_path)
                        except:
                            pass
            
            # 如果转换失败，使用原始提取方法
            return EnhancedFileAnalyzer._extract_doc_content(file_path)
            
        except Exception as e:
            return f".doc文件解析错误: {str(e)}"
        
    @staticmethod
    def _extract_doc_content(file_path):
        """提取旧版Word文档内容"""
        try:
            # 方法1: 尝试使用antiword（需要安装）
            try:
                import subprocess
                result = subprocess.run(['antiword', file_path], capture_output=True, text=True, timeout=30)
                if result.returncode == 0:
                    return result.stdout
            except:
                pass
            
            # 方法2: 尝试使用catdoc（需要安装）
            try:
                import subprocess
                result = subprocess.run(['catdoc', file_path], capture_output=True, text=True, timeout=30)
                if result.returncode == 0:
                    return result.stdout
            except:
                pass
            
            # 方法3: 使用python-docx兼容模式
            try:
                import docx
                # 尝试以二进制方式读取
                with open(file_path, 'rb') as f:
                    content = f.read().decode('utf-8', errors='ignore')
                    return content[:5000]  # 限制长度
            except:
                pass
            
            return "需要安装antiword或catdoc来解析.doc文件"
            
        except Exception as e:
            return f".doc文件解析错误: {str(e)}"
    
    @staticmethod
    def _extract_pdf_content(file_path):
        """提取PDF内容"""
        try:
            # 优先使用PyPDF2
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    return text if text else "PDF内容为空或受保护"
            except ImportError:
                return "请安装PyPDF2: pip install PyPDF2"
            except Exception as e:
                return f"PDF解析错误: {str(e)}"
                
        except Exception as e:
            return f"PDF读取失败: {str(e)}"
    
    @staticmethod
    def _extract_excel_content(file_path):
        """提取Excel内容"""
        try:
            import pandas as pd
            try:
                # 读取所有工作表
                excel_file = pd.ExcelFile(file_path)
                content = []
                
                for sheet_name in excel_file.sheet_names:
                    # 只读取前5行作为样本
                    df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=5)
                    content.append(f"工作表: {sheet_name}")
                    content.append(str(df))
                    
                return '\n'.join(content) if content else "Excel文件为空"
                
            except Exception as e:
                return f"Excel读取错误: {str(e)}"
                
        except ImportError:
            return "请安装pandas: pip install pandas openpyxl"
    
    @staticmethod
    def _extract_ppt_content(file_path):
        """提取PPT内容"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                text_content.append(f"--- 幻灯片 {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
            
            return '\n'.join(text_content) if text_content else "PPT内容为空"
            
        except ImportError:
            return "请安装python-pptx: pip install python-pptx"
        except Exception as e:
            return f"PPT读取错误: {str(e)}"

class FileInfoExtractor:
    """文件信息提取器"""
    
    @staticmethod
    def extract_file_info(file_path):
        """提取文件基本信息"""
        try:
            file_stat = os.stat(file_path)
            file_name = os.path.basename(file_path)
            
            return {
                'name': file_name,
                'path': os.path.abspath(file_path),
                'size': file_stat.st_size,
            
                'update_time': datetime.fromtimestamp(file_stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
                'extension': os.path.splitext(file_name)[1],
                'directory': os.path.dirname(os.path.abspath(file_path)),
                'hash': hashlib.md5(open(file_path, 'rb').read()).hexdigest()[:16]
            }
        except Exception as e:
            logger.error(f"文件信息提取失败: {str(e)}")
            return None

class EnhancedKnowledgeBaseUploader:
    """增强的知识库上传器 - 简化版本"""
    
    def __init__(self):
        # 初始化所有必要属性
        self.base_url = config.DIFY_BASE_URL
        self.api_key = config.DATASET_API_KEY
        self.converter = FileConverter()
        self.failed_conversions = set()  # 记录转换失败的文件
    
    def _get_mime_type(self, file_ext):
        """获取文件的MIME类型"""
        mime_types = {
            '.txt': 'text/plain',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.pdf': 'application/pdf',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.csv': 'text/csv',
            '.md': 'text/markdown',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.bmp': 'image/bmp',
            '.webp': 'image/webp',
            '.tiff': 'image/tiff'
        }
        return mime_types.get(file_ext, 'application/octet-stream')
    
    def upload_file(self, file_path, knowledge_base_id=None, use_parent_child_mode=False):
        """简化的文件上传方法 - 让Dify使用默认设置"""
        temp_file_created = False
        temp_file_path = None
        
        try:
            if not self.api_key:
                logger.warning("知识库API密钥未设置，跳过上传")
                return False
            
            if knowledge_base_id is None:
                knowledge_base_id = config.TXT_KNOWLEDGE_BASE_ID
            
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1].lower()
            
            # 检查是否需要转换格式
            upload_path = file_path
            upload_name = file_name
            
            # 检查是否已经尝试过转换但失败
            file_id = f"{file_path}_{os.path.getsize(file_path)}"
            if file_id in self.failed_conversions:
                logger.info(f"跳过已失败的文件转换: {file_name}")
            elif file_ext == '.doc':
                logger.info(f"检测到.doc文件，尝试转换为.docx: {file_name}")
                converted_path = self.converter.convert_doc_to_docx(file_path)
                if converted_path and os.path.exists(converted_path):
                    upload_path = converted_path
                    upload_name = os.path.splitext(file_name)[0] + '.docx'
                    temp_file_created = True
                    temp_file_path = converted_path
                    logger.info(f"使用转换后的文件进行上传: {file_name} -> {upload_name}")
                else:
                    logger.warning(f"文件转换失败，记录失败状态: {file_name}")
                    self.failed_conversions.add(file_id)
            
            url = f"{config.DIFY_BASE_URL}/v1/datasets/{knowledge_base_id}/document/create-by-file"
            
            # 检查文件大小
            file_size = os.path.getsize(upload_path)
            if file_size > 100 * 1024 * 1024:  # 100MB限制
                logger.warning(f"文件过大({file_size}字节)，跳过上传: {upload_name}")
                return False
            
            # 读取文件内容
            with open(upload_path, 'rb') as file:
                file_content = file.read()
            
            # 使用BytesIO避免文件锁定
            from io import BytesIO
            file_stream = BytesIO(file_content)
            
            upload_ext = os.path.splitext(upload_name)[1].lower()
            mime_type = self._get_mime_type(upload_ext)
            
            files = {'file': (upload_name, file_stream, mime_type)}
            
            # 🔧🔧 关键修改：使用最简单的配置，只提供文件名
            # 不指定任何处理规则，让Dify完全使用知识库的默认设置
            data = {
                'name': upload_name
                # 不指定 process_rule，让Dify自动处理
                # 不指定 indexing_technique，使用知识库默认设置
            }
            
            if use_parent_child_mode:
                logger.info(f"上传到父子模式知识库（使用Dify默认设置）: {upload_name}")
            else:
                logger.info(f"上传到普通知识库: {upload_name}")
            
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "User-Agent": "FileMonitor/1.0"
            }
            
            logger.debug(f"上传请求数据 - 文件名: {upload_name}, 知识库ID: {knowledge_base_id}")
            
            response = requests.post(
                url, 
                headers=headers, 
                files=files, 
                data=data, 
                timeout=config.API_TIMEOUT
            )
            
            if response.status_code in [200, 201]:
                kb_type = "父子模式知识库" if use_parent_child_mode else "知识库"
                logger.info(f"✅ {kb_type}上传成功: {upload_name} -> 知识库 {knowledge_base_id}")
                
                # 打印成功响应信息
                try:
                    result = response.json()
                    logger.debug(f"上传成功响应: {json.dumps(result, ensure_ascii=False)[:200]}...")
                except:
                    pass
                    
                return True
            else:
                error_msg = f"❌❌ 知识库上传失败: {response.status_code} - {response.text}"
                logger.error(error_msg)
                
                # 详细错误分析
                error_text = response.text.lower()
                if "doc_form" in error_text or "segmentation" in error_text:
                    logger.error("📋📋 文档格式错误：可能与知识库的分段模式设置有关")
                elif "not found" in error_text:
                    logger.error("🔍🔍 知识库不存在或API密钥无权限")
                elif "indexing_technique" in error_text:
                    logger.error("⚙⚙️ 索引技术错误")
                elif "unauthorized" in error_text:
                    logger.error("🔐🔐 认证失败：请检查API密钥")
                    
                return False
                
        except requests.exceptions.Timeout:
            logger.error("⏰⏰⏰ 上传请求超时")
            return False
        except requests.exceptions.ConnectionError:
            logger.error("🌐🌐 网络连接错误，请检查Dify服务是否可用")
            return False
        except Exception as e:
            logger.error(f"💥💥 上传过程异常: {str(e)}")
            return False
        finally:
            # 清理临时文件
            if temp_file_created and temp_file_path and os.path.exists(temp_file_path):
                self._safe_delete_file(temp_file_path)
    
    def _safe_delete_file(self, file_path):
        """安全删除文件"""
        import time
        max_retries = 5
        for attempt in range(max_retries):
            try:
                os.remove(file_path)
                logger.info(f"已清理临时文件: {os.path.basename(file_path)}")
                return True
            except PermissionError:
                if attempt < max_retries - 1:
                    time.sleep(1)
                    continue
                else:
                    logger.warning(f"无法删除临时文件(重试{max_retries}次后): {os.path.basename(file_path)}")
                    return False
            except Exception as e:
                logger.warning(f"删除临时文件失败: {str(e)}")
                return False
            
class FileMonitor:
    """文件监控器"""
    
    def __init__(self):
        self.index_generator = DifyChatflowIndexGenerator()
        self.uploader = EnhancedKnowledgeBaseUploader()
        self.processed_files = {}  # 文件ID -> 处理时间
        self.processing_lock = set()  # 正在处理的文件锁
    
    def should_process(self, file_path):
        """判断是否应该处理文件"""
        file_name = os.path.basename(file_path)
        
        # 跳过系统文件和标记文件
        skip_markers = ['_summary', '_index', '_workflow', '_dify', '_chatflow', '~$', '_converted']
        if any(marker in file_name for marker in skip_markers):
            return False
        
        # 检查文件扩展名
        file_ext = os.path.splitext(file_name)[1].lower()
        if file_ext not in config.ALLOWED_EXTENSIONS:
            return False
        
        # 生成文件唯一标识（路径+大小+修改时间）
        try:
            stat = os.stat(file_path)
            file_id = f"{file_path}_{stat.st_size}_{stat.st_mtime}"
        except:
            file_id = file_path
        
        # 检查是否正在处理
        if file_id in self.processing_lock:
            return False
        
        # 检查是否最近处理过
        current_time = time.time()
        if file_id in self.processed_files:
            time_diff = current_time - self.processed_files[file_id]
            if time_diff < config.PROCESS_INTERVAL:
                return False
        
        return True
    
    def process_file(self, file_path):
        """处理文件 - 图片跳过原文件上传"""
        file_id = None
        try:
            if not self.should_process(file_path):
                return
            
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1].lower()
            is_image = file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff']
            
            logger.info(f"开始处理文件: {file_name} (类型: {'图片' if is_image else '文档'})")
            
            # 生成文件ID并加锁
            try:
                stat = os.stat(file_path)
                file_id = f"{file_path}_{stat.st_size}_{stat.st_mtime}"
            except:
                file_id = file_path
            
            self.processing_lock.add(file_id)
            
            # 等待文件稳定
            time.sleep(2)
            
            # 1. 生成索引文件（图片和文档都执行）
            index_path = self.index_generator.generate_index_file(file_path)
            
            if index_path:
                # 2. 上传索引文件到.txt知识库
                index_success = self.uploader.upload_file(index_path, config.TXT_KNOWLEDGE_BASE_ID)
                if index_success:
                    logger.info(f"索引文件上传成功: {file_name}")
            
            # 3. 只有非图片文件才上传原文件到知识库
            if not is_image:
                original_kb_id = config.ACTUAL_ORIGINAL_KB_ID
                use_parent_child_mode = config.PARENT_CHILD_KB_ENABLED and config.PARENT_CHILD_KB_ID
                
                if use_parent_child_mode:
                    kb_type = "父子模式知识库"
                    original_success = self.uploader.upload_file(
                        file_path, 
                        original_kb_id, 
                        use_parent_child_mode=True
                    )
                else:
                    kb_type = "原文件库"
                    original_success = self.uploader.upload_file(file_path, original_kb_id)
                
                if original_success:
                    logger.info(f"原文件上传成功到{kb_type}: {file_name} -> {original_kb_id}")
                else:
                    logger.error(f"原文件上传到{kb_type}失败: {file_name}")
            else:
                logger.info(f"图片文件跳过原文件上传: {file_name}")
            
            # 记录处理时间
            self.processed_files[file_id] = time.time()
            
        except Exception as e:
            logger.error(f"处理文件异常: {file_name} - {str(e)}")
        finally:
            # 释放处理锁
            if file_id and file_id in self.processing_lock:
                self.processing_lock.remove(file_id)
    
    def open_image_by_filename(self, filename):
        """根据文件名打开图片"""
        try:
            image_path = os.path.join(config.MONITOR_DIR, filename)
            if os.path.exists(image_path):
                # 使用系统默认程序打开图片
                if sys.platform == "win32":
                    os.startfile(image_path)
                elif sys.platform == "darwin":  # macOS
                    import subprocess
                    subprocess.run(["open", image_path])
                else:  # Linux
                    import subprocess
                    subprocess.run(["xdg-open", image_path])
                logger.info(f"已打开图片: {filename}")
                return True
            else:
                logger.error(f"图片文件不存在: {filename}")
                return False
        except Exception as e:
            logger.error(f"打开图片失败: {filename} - {str(e)}")
            return False

class FileEventHandler(FileSystemEventHandler):
    """文件事件处理器"""
    
    def __init__(self, monitor):
        self.monitor = monitor
    
    def on_created(self, event):
        if not event.is_directory:
            self.monitor.process_file(event.src_path)
    
    def on_modified(self, event):
        if not event.is_directory:
            self.monitor.process_file(event.src_path)

def start_monitoring():
    """启动文件监控"""
    try:
        # 检查监控目录
        if not os.path.exists(config.MONITOR_DIR):
            os.makedirs(config.MONITOR_DIR)
            logger.info(f"创建监控目录: {config.MONITOR_DIR}")
        
        # 显示配置信息
        logger.info("=== Dify Chatflow文件分析系统 ===")
        logger.info(f"监控目录: {config.MONITOR_DIR}")
        logger.info(f"支持格式: {', '.join(config.ALLOWED_EXTENSIONS)}")
        
        # 显示知识库配置
        if config.PARENT_CHILD_KB_ENABLED and config.PARENT_CHILD_KB_ID:
            logger.info(f"原文件库: 父子模式知识库 -> {config.PARENT_CHILD_KB_ID}")
            logger.info("原文件处理模式: 使用Dify默认设置")
        else:
            logger.info(f"原文件库: 普通知识库 -> {config.ORIGINAL_KNOWLEDGE_BASE_ID}")
            logger.info("原文件处理模式: 普通模式")
        
        logger.info(f"索引文件库: {config.TXT_KNOWLEDGE_BASE_ID}")
        logger.info(f"Chatflow APP ID: {config.CHATFLOW_APP_ID}")
        
        # 创建监控器
        monitor = FileMonitor()
        event_handler = FileEventHandler(monitor)
        observer = Observer()
        observer.schedule(event_handler, path=config.MONITOR_DIR, recursive=False)
        observer.start()
        
        logger.info("✅ Dify Chatflow监控服务已启动")
        logger.info("📁 将文件放入监控目录，系统将自动处理")
        
        # 保持运行
        try:
            while True:
                time.sleep(1)
        except KeyboardInterrupt:
            logger.info("监控服务已停止")
            observer.stop()
        observer.join()
            
    except Exception as e:
        logger.error(f"监控服务启动失败: {str(e)}")

# 在文件末尾添加测试函数
def test_multimodal_support():
    """测试多模态支持"""
    analyzer = DifyChatflowAnalyzer()
    test_multimodal_support()
    # 检查支持情况
    support = analyzer.check_multimodal_support()
    
    if support:
        print("✅ 您的Dify环境支持多模态分析")
        print("需要修改API调用方式")
    else:
        print("❌ 您的Dify环境不支持多模态分析")
        print("建议使用备用方案：基于图片元数据生成描述")

if __name__ == "__main__":
    print("=== Dify Chatflow文件分析系统 ===")
    print(f"监控目录: {config.MONITOR_DIR}")
    print("分析模式: 优先使用Dify Chatflow，失败时使用本地规则推断")
    print("上传策略: 原文件 -> 父子模式知识库, 索引文件 -> .txt文件库")
    print("图片处理: 生成索引文件，跳过原文件上传")
    print("正在启动...")
    
    # 检查必要的库
    try:
        import docx
        print("✅ python-docx 已安装")
    except ImportError:
        print("⚠️ python-docx 未安装，Word文档解析将受限")
    
    try:
        import PyPDF2
        print("✅ PyPDF2 已安装")
    except ImportError:
        print("⚠️ PyPDF2 未安装，PDF文档解析将受限")
    
    if PIL_AVAILABLE:
        print("✅ Pillow 已安装，支持图片分析")
    else:
        print("⚠️ Pillow 未安装，请运行: pip install Pillow")
    
    start_monitoring()